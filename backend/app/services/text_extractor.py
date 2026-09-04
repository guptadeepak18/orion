import io
import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)

_TEXT_CACHE: dict = {}


def extract_text_from_file(file_path: str, filename: Optional[str] = None, file_bytes: Optional[bytes] = None) -> str:
    """
    Extract clean textual content from various document formats (PDF, DOCX, PPTX, XLSX, TXT)
    for LLM evaluation, supporting local paths, R2 storage keys, and raw byte buffers.
    """
    if file_bytes is None and file_path and file_path in _TEXT_CACHE:
        return _TEXT_CACHE[file_path]

    name = (filename or os.path.basename(file_path)).lower()
    ext = os.path.splitext(name)[1].lower()

    if file_bytes is None:
        try:
            from app.services.storage_service import storage_service
            file_bytes = storage_service.read_file_bytes_sync(file_path)
        except Exception as e:
            logger.warning(f"[TextExtractor] Failed to read bytes for '{file_path}': {e}")
            return ""

    if not file_bytes:
        return ""

    try:
        if ext == ".pdf":
            res = _extract_pdf_from_bytes(file_bytes)
        elif ext in [".docx", ".doc"]:
            res = _extract_docx_from_bytes(file_bytes)
        elif ext in [".pptx", ".ppt"]:
            res = _extract_pptx_from_bytes(file_bytes)
        elif ext in [".xlsx", ".xls", ".csv"]:
            res = _extract_spreadsheet_from_bytes(file_bytes, ext)
        elif ext in [".txt", ".md", ".json", ".rtf", ".py", ".html"]:
            res = _extract_text_from_bytes(file_bytes)
        else:
            res = _extract_text_from_bytes(file_bytes)

        if file_path and res:
            _TEXT_CACHE[file_path] = res
        return res
    except Exception as e:
        logger.error(f"Error extracting text from {name}: {e}")
        return f"[Error extracting text from {name}: {str(e)}]"


def _extract_pdf_from_bytes(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    text_parts = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if page_text.strip():
            text_parts.append(f"--- Page {i+1} ---\n" + page_text.strip())
    return "\n\n".join(text_parts)


def _extract_docx_from_bytes(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    lines = []
    for p in doc.paragraphs:
        if p.text.strip():
            lines.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                lines.append(f"| {row_text} |")
    return "\n".join(lines)


def _extract_pptx_from_bytes(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_lines))
    return "\n\n".join(slides_text)


def _extract_spreadsheet_from_bytes(data: bytes, ext: str) -> str:
    if ext == ".csv":
        # Try multiple encodings for CSV
        for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                return data.decode(encoding)
            except (UnicodeDecodeError, ValueError):
                continue
        return data.decode("utf-8", errors="ignore")

    if ext == ".xls":
        return _extract_xls_from_bytes(data)

    # .xlsx — use openpyxl
    return _extract_xlsx_from_bytes(data)


def _extract_xlsx_from_bytes(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows_text = []
        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            row_count += 1
            if row_count > 5000:
                rows_text.append(f"... (truncated after {row_count} rows)")
                break
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheets_text.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows_text))
    wb.close()
    return "\n\n".join(sheets_text) if sheets_text else "[Excel file contained no readable cell data]"


def _extract_xls_from_bytes(data: bytes) -> str:
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=data)
        sheets_text = []
        for sheet_name in wb.sheet_names():
            sheet = wb.sheet_by_name(sheet_name)
            rows_text = []
            max_rows = min(sheet.nrows, 5000)
            for rx in range(max_rows):
                cells = []
                for cx in range(sheet.ncols):
                    cell = sheet.cell(rx, cx)
                    val = cell.value
                    if val is None or (isinstance(val, str) and not val.strip()):
                        continue
                    # Format dates properly
                    if cell.ctype == xlrd.XL_CELL_DATE:
                        try:
                            dt_tuple = xlrd.xldate_as_tuple(val, wb.datemode)
                            val = f"{dt_tuple[0]:04d}-{dt_tuple[1]:02d}-{dt_tuple[2]:02d}"
                        except Exception:
                            val = str(val)
                    elif cell.ctype == xlrd.XL_CELL_NUMBER:
                        # Display integers without decimals
                        val = int(val) if val == int(val) else val
                    cells.append(str(val).strip())
                if cells:
                    rows_text.append(" | ".join(cells))
            if sheet.nrows > 5000:
                rows_text.append(f"... (truncated after 5000 of {sheet.nrows} rows)")
            if rows_text:
                sheets_text.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows_text))
        return "\n\n".join(sheets_text) if sheets_text else "[XLS file contained no readable cell data]"
    except ImportError:
        logger.warning("[TextExtractor] xlrd not installed; cannot read .xls files. Install with: pip install xlrd")
        return "[Error: .xls file format requires xlrd library which is not installed]"
    except Exception as e:
        logger.error(f"[TextExtractor] Failed to read .xls file: {e}")
        return f"[Error reading .xls file: {str(e)}]"


def _extract_text_from_bytes(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")
